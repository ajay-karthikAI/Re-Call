from datetime import datetime
from pathlib import Path
from uuid import UUID

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy.ext.asyncio import AsyncSession

from models import Meeting
from services.chart_export_service import chart_cards_for_export, code_snippets_for_export, display_value, numeric_rows
from services.export_formatting import format_duration


NAVY = RGBColor(30, 39, 97)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(20, 27, 36)
MUTED = RGBColor(93, 104, 121)
GREEN = RGBColor(0, 229, 160)
CODE_BG = RGBColor(13, 17, 23)
LIGHT = RGBColor(244, 247, 250)
EXPORT_FONT = "Times New Roman"


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, font=EXPORT_FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _title(slide, text: str, dark: bool = False) -> None:
    _textbox(
        slide,
        Inches(0.7),
        Inches(0.45),
        Inches(8.6),
        Inches(0.55),
        text,
        size=32 if not dark else 40,
        color=WHITE if dark else INK,
        bold=True,
    )
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.1), Inches(1.2), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()


def _add_bullets(slide, bullets: list[str], x, y, w, h, color=INK) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = EXPORT_FONT
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)


def _style_cell_text(cell, color=INK, bold=False, size=12) -> None:
    cell.text_frame.word_wrap = True
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = EXPORT_FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        for run in paragraph.runs:
            run.font.name = EXPORT_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold


def _summary_bullets(summary: str) -> list[str]:
    parts = [item.strip() for item in summary.replace("\n", " ").split(".") if item.strip()]
    return [f"{item}." for item in parts[:5]] or ["No summary was generated."]


def _notes(meeting: Meeting) -> dict:
    return meeting.notes_json or {}


def _next_step_text(step) -> str:
    if not isinstance(step, dict):
        return str(step)
    priority = str(step.get("priority") or "next").upper()
    task = str(step.get("task") or "Follow up")
    reason = str(step.get("reason") or "").strip()
    return f"[{priority}] {task}" + (f" - {reason}" if reason else "")


def _add_chart_slides(prs: Presentation, blank, charts: list[dict]) -> None:
    for chart in charts:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, WHITE)
        _title(slide, str(chart.get("title") or "Chart")[:64])
        chart_type = chart.get("chart_type") or "bar_chart"
        if chart_type == "needs_data":
            _draw_missing_chart(slide, chart)
        elif chart_type == "timeline":
            _draw_timeline(slide, chart)
        elif chart_type == "table":
            _draw_table_chart(slide, chart)
        elif chart_type == "line_chart":
            _draw_line_chart(slide, chart)
        else:
            _draw_bar_chart(slide, chart)
        if chart.get("insight"):
            _textbox(slide, Inches(0.9), Inches(6.65), Inches(11.4), Inches(0.28), str(chart["insight"])[:180], 12, MUTED)


def _draw_missing_chart(slide, chart: dict) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.65), Inches(10.9), Inches(1.35))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(224, 230, 237)
    _textbox(slide, Inches(1.25), Inches(1.98), Inches(10.1), Inches(0.28), "Graph data missing", 19, INK, True)
    prompt = chart.get("missing_data_prompt") or "A graph request was detected, but structured chart data was not saved."
    _textbox(slide, Inches(1.25), Inches(2.42), Inches(10.0), Inches(0.32), str(prompt)[:180], 13, MUTED)


def _ppt_tick_values(min_value: float, max_value: float) -> list[float]:
    low = min(min_value, 0)
    high = max(max_value, 1)
    if low == high:
        high = low + 1
    return [low, low + (high - low) / 2, high]


def _ppt_axis_value(value: float) -> str:
    if abs(value) >= 1000:
        return f"{value:,.0f}"
    return str(int(value)) if float(value).is_integer() else f"{value:.1f}"


def _draw_bar_chart(slide, chart: dict) -> None:
    rows = numeric_rows(chart)
    if not rows:
        _textbox(slide, Inches(0.95), Inches(1.65), Inches(10.5), Inches(0.35), "No chart data available.", 18, MUTED)
        return
    max_value = max(abs(float(row.get("value") or 0)) for row in rows) or 1
    x_label = Inches(0.95)
    x_bar = Inches(3.7)
    y = Inches(1.55)
    max_width = Inches(7.2)
    plot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(1.28), Inches(7.75), Inches(5.1))
    plot.fill.solid()
    plot.fill.fore_color.rgb = RGBColor(247, 250, 246)
    plot.line.color.rgb = RGBColor(216, 226, 214)
    for tick in _ppt_tick_values(0, max_value):
        x = x_bar + max_width * (tick / max_value)
        grid = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, Inches(1.42), x, Inches(6.22))
        grid.line.color.rgb = RGBColor(224, 232, 222)
        _textbox(slide, x - Inches(0.2), Inches(6.32), Inches(0.55), Inches(0.18), _ppt_axis_value(tick), 8, MUTED)
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_bar, Inches(6.18), x_bar + max_width, Inches(6.18))
    axis.line.color.rgb = MUTED
    if chart.get("y_label"):
        _textbox(slide, x_label, Inches(1.28), Inches(2.4), Inches(0.2), str(chart["y_label"])[:28], 9, MUTED, True)
    for row in rows[:8]:
        label = str(row.get("label") or "")[:34]
        _textbox(slide, x_label, y + Inches(0.04), Inches(2.45), Inches(0.24), label, 11, INK, True)
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x_bar, y, max_width, Inches(0.24))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT
        bg.line.fill.background()
        width = max(Inches(0.12), max_width * (abs(float(row.get("value") or 0)) / max_value))
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x_bar, y, width, Inches(0.24))
        bar.fill.solid()
        bar.fill.fore_color.rgb = GREEN
        bar.line.fill.background()
        label_x = min(x_bar + width + Inches(0.12), x_bar + max_width - Inches(0.62))
        _textbox(slide, label_x, y + Inches(0.02), Inches(0.9), Inches(0.22), display_value(row), 11, INK, True)
        y += Inches(0.56)
    if chart.get("x_label"):
        _textbox(slide, x_bar + max_width - Inches(0.65), Inches(6.62), Inches(1.1), Inches(0.18), str(chart["x_label"])[:22], 8, MUTED)


def _draw_line_chart(slide, chart: dict) -> None:
    rows = numeric_rows(chart)
    if not rows:
        _textbox(slide, Inches(0.95), Inches(1.65), Inches(10.5), Inches(0.35), "No chart data available.", 18, MUTED)
        return
    values = [float(row.get("value") or 0) for row in rows]
    min_value = min(min(values), 0)
    max_value = max(values) or 1
    value_range = max(max_value - min_value, 1)
    x0 = Inches(0.95)
    y0 = Inches(5.85)
    width = Inches(11.2)
    height = Inches(4.15)

    plot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(1.35), Inches(11.75), Inches(4.85))
    plot.fill.solid()
    plot.fill.fore_color.rgb = RGBColor(247, 250, 246)
    plot.line.color.rgb = RGBColor(216, 226, 214)

    for tick in _ppt_tick_values(min_value, max_value):
        y = y0 - height * ((tick - min_value) / value_range)
        grid = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x0, y, x0 + width, y)
        grid.line.color.rgb = RGBColor(224, 232, 222)
        _textbox(slide, Inches(0.42), y - Inches(0.08), Inches(0.45), Inches(0.16), _ppt_axis_value(tick), 8, MUTED)

    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x0, y0, x0 + width, y0)
    axis.line.color.rgb = MUTED
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x0, y0, x0, y0 - height)
    y_axis.line.color.rgb = MUTED

    points = []
    for index, row in enumerate(rows):
        x = x0 + (width / 2 if len(rows) == 1 else width * (index / (len(rows) - 1)))
        y = y0 - height * ((float(row.get("value") or 0) - min_value) / value_range)
        points.append((x, y, row))
    for start, end in zip(points, points[1:]):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start[0], start[1], end[0], end[1])
        connector.line.color.rgb = GREEN
        connector.line.width = Pt(2.5)
    for x, y, row in points:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - Inches(0.06), y - Inches(0.06), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()
        _textbox(slide, x - Inches(0.28), max(Inches(1.23), y - Inches(0.3)), Inches(0.7), Inches(0.18), display_value(row), 8, INK, True)
        _textbox(slide, x - Inches(0.34), y0 + Inches(0.12), Inches(0.75), Inches(0.2), str(row.get("label") or "")[:10], 8, MUTED)
    if chart.get("y_label"):
        _textbox(slide, Inches(0.95), Inches(1.1), Inches(3.0), Inches(0.22), str(chart["y_label"])[:28], 10, MUTED, True)
    if chart.get("x_label"):
        _textbox(slide, Inches(11.3), Inches(6.1), Inches(0.9), Inches(0.18), str(chart["x_label"])[:20], 8, MUTED)


def _draw_timeline(slide, chart: dict) -> None:
    rows = chart.get("data") or []
    if not rows:
        _textbox(slide, Inches(0.95), Inches(1.65), Inches(10.5), Inches(0.35), "No timeline data available.", 18, MUTED)
        return
    x = Inches(1.15)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.55), Inches(0.035), Inches(4.72))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()
    y = Inches(1.48)
    for row in rows[:6]:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - Inches(0.09), y + Inches(0.04), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()
        _textbox(slide, Inches(1.55), y, Inches(2.1), Inches(0.28), str(row.get("label") or "")[:26], 15, INK, True)
        _textbox(slide, Inches(3.55), y + Inches(0.02), Inches(8.3), Inches(0.35), str(row.get("text") or display_value(row))[:140], 13, MUTED)
        y += Inches(0.78)


def _draw_table_chart(slide, chart: dict) -> None:
    rows = chart.get("data") or []
    if not rows:
        _textbox(slide, Inches(0.95), Inches(1.65), Inches(10.5), Inches(0.35), "No table data available.", 18, MUTED)
        return
    table_shape = slide.shapes.add_table(len(rows[:8]) + 1, 2, Inches(0.95), Inches(1.55), Inches(11.3), Inches(4.7))
    table = table_shape.table
    table.columns[0].width = Inches(5.4)
    table.columns[1].width = Inches(5.9)
    for col, header in enumerate(["Label", "Value"]):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _style_cell_text(cell, WHITE, True)
    for row_index, row in enumerate(rows[:8], start=1):
        for col, value in enumerate([row.get("label"), row.get("text") or display_value(row)]):
            cell = table.cell(row_index, col)
            cell.text = str(value or "")
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else LIGHT
            _style_cell_text(cell)


async def generate(session: AsyncSession, meeting_id: UUID) -> Path:
    meeting = await session.get(Meeting, meeting_id)
    if meeting is None:
        raise ValueError(f"Meeting {meeting_id} was not found")

    notes = _notes(meeting)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, NAVY)
    _textbox(slide, Inches(0.8), Inches(1.6), Inches(10.8), Inches(0.9), meeting.title, 40, WHITE, True)
    _textbox(
        slide,
        Inches(0.85),
        Inches(2.62),
        Inches(8.5),
        Inches(0.4),
        f"{meeting.created_at:%B %d, %Y} | {format_duration(meeting.duration_seconds)}",
        18,
        GREEN,
    )
    mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(4.55), Inches(2.2), Inches(0.56))
    mark.fill.solid()
    mark.fill.fore_color.rgb = GREEN
    mark.line.fill.background()
    _textbox(slide, Inches(1.05), Inches(4.68), Inches(1.8), Inches(0.25), "Re: Call", 16, NAVY, True)

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Participants")
    participants = notes.get("participants") or ["Speaker 1"]
    for index, name in enumerate(participants[:10]):
        col = index % 2
        row = index // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.55 + row * 0.9)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(5.2), Inches(0.58))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = RGBColor(224, 230, 237)
        initials = "".join(part[0] for part in str(name).split()[:2]).upper() or "S"
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.16), y + Inches(0.11), Inches(0.36), Inches(0.36))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN
        circle.line.fill.background()
        _textbox(slide, x + Inches(0.235), y + Inches(0.19), Inches(0.24), Inches(0.16), initials[:2], 8, NAVY, True)
        _textbox(slide, x + Inches(0.68), y + Inches(0.17), Inches(4.0), Inches(0.25), str(name), 16, INK, True)

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Executive Summary")
    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.52), Inches(11.7), Inches(0.28))
    quote.fill.solid()
    quote.fill.fore_color.rgb = GREEN
    quote.line.fill.background()
    _add_bullets(slide, _summary_bullets(notes.get("summary", "")), Inches(1.0), Inches(2.0), Inches(10.8), Inches(3.6))

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Insights")
    insights = notes.get("insights") or ["No insights were captured."]
    _add_bullets(slide, [str(item) for item in insights[:7]], Inches(1.0), Inches(1.65), Inches(10.8), Inches(4.5))

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Key Decisions")
    decisions = notes.get("key_decisions") or ["No key decisions were captured."]
    for index, decision in enumerate(decisions[:6]):
        y = Inches(1.55 + index * 0.78)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.9), y, Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN
        circle.line.fill.background()
        _textbox(slide, Inches(1.0), y + Inches(0.09), Inches(0.28), Inches(0.18), "OK", 8, NAVY, True)
        _textbox(slide, Inches(1.55), y + Inches(0.03), Inches(10.4), Inches(0.36), str(decision), 17, INK)

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Action Items")
    actions = notes.get("action_items") or [{"owner": "TBD", "task": "No action items were captured.", "due": "TBD"}]
    table_shape = slide.shapes.add_table(len(actions[:8]) + 1, 3, Inches(0.8), Inches(1.55), Inches(11.7), Inches(4.6))
    table = table_shape.table
    widths = [2.1, 7.2, 2.4]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for col, header in enumerate(["Owner", "Task", "Due Date"]):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _style_cell_text(cell, WHITE, True)
    for row, action in enumerate(actions[:8], start=1):
        for col, value in enumerate([action.get("owner", "TBD"), action.get("task", ""), action.get("due", "TBD")]):
            cell = table.cell(row, col)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row % 2 else LIGHT
            _style_cell_text(cell)

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Next Steps")
    next_steps = notes.get("next_steps") or ["No next steps were captured."]
    _add_bullets(slide, [_next_step_text(step) for step in next_steps[:7]], Inches(1.0), Inches(1.65), Inches(10.8), Inches(4.5))

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, WHITE)
    _title(slide, "Topics Covered")
    topics = notes.get("topics_discussed") or ["No topics were captured."]
    x = Inches(0.85)
    y = Inches(1.65)
    for topic in topics[:18]:
        label = str(topic)[:34]
        width = Inches(max(1.6, min(4.0, 0.18 * len(label) + 0.7)))
        if x + width > Inches(12.1):
            x = Inches(0.85)
            y += Inches(0.72)
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, width, Inches(0.42))
        pill.fill.solid()
        pill.fill.fore_color.rgb = GREEN
        pill.line.fill.background()
        _textbox(slide, x + Inches(0.18), y + Inches(0.12), width - Inches(0.32), Inches(0.16), label, 12, NAVY, True)
        x += width + Inches(0.25)

    charts = chart_cards_for_export(notes, meeting.transcript or "")
    _add_chart_slides(prs, blank, charts)

    snippets = code_snippets_for_export(notes, charts)
    for snippet in snippets:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, CODE_BG)
        _title(slide, snippet.get("description", "Code"), dark=True)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(1.35), Inches(1.5), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        _textbox(slide, Inches(0.96), Inches(1.45), Inches(1.1), Inches(0.12), snippet.get("language", "code"), 10, CODE_BG, True)
        code = str(snippet.get("code", ""))[:2400]
        _textbox(slide, Inches(0.85), Inches(1.9), Inches(11.8), Inches(4.85), code, 11, WHITE)

    slide = prs.slides.add_slide(blank)
    _set_bg(slide, NAVY)
    _textbox(slide, Inches(0.85), Inches(2.4), Inches(8.8), Inches(0.7), "Notes by Re: Call", 40, WHITE, True)
    _textbox(slide, Inches(0.9), Inches(3.25), Inches(5.0), Inches(0.35), datetime.now().strftime("%B %d, %Y"), 18, GREEN)
    footer_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.15), Inches(3.4), Inches(0.58))
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = GREEN
    footer_shape.line.fill.background()
    _textbox(slide, Inches(1.13), Inches(5.31), Inches(2.9), Inches(0.18), "Generated after the meeting", 12, NAVY, True)

    output = Path("/tmp") / f"{meeting_id}.pptx"
    prs.save(output)
    return output
